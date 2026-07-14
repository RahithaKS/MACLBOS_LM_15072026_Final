import sys
sys.path.insert(0, '/home/runner/workspace/.pythonlibs/lib/python3.11/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)
C_NAVY   = RGBColor(0x1A, 0x35, 0x5E)
C_BLUE   = RGBColor(0x1E, 0x6F, 0xBF)
C_ACCENT = RGBColor(0x00, 0xC2, 0xCB)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xEC, 0xF0, 0xF1)
C_MID    = RGBColor(0xBD, 0xC3, 0xC7)
C_YELLOW = RGBColor(0xF3, 0x9C, 0x12)
C_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
C_PINK   = RGBColor(0xE9, 0x1E, 0x63)
C_DKRED  = RGBColor(0x2C, 0x10, 0x10)
C_DKGRN  = RGBColor(0x0D, 0x2A, 0x1A)

def R(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=10, bold=False, clr=C_WHITE, al=PP_ALIGN.LEFT, it=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = al
        r = p.add_run(); r.text = ln
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = it; r.font.color.rgb = clr

def bg(slide, c=C_DARK):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def header(slide, title, sub=None):
    add_rect(slide, 0, 0, 13.33, 1.28, fill=C_NAVY)
    add_rect(slide, 0, 1.28, 13.33, 0.06, fill=C_BLUE)
    txt(slide, title, 0.38, 0.1, 12.5, 0.72, sz=23, bold=True)
    if sub: txt(slide, sub, 0.38, 0.8, 12.5, 0.42, sz=11, clr=C_ACCENT, it=True)

def badge(slide, l, t, w, h, text, c=C_BLUE):
    add_rect(slide, l, t, w, h, fill=c)
    txt(slide, text, l+0.05, t+0.03, w-0.1, h-0.06, sz=8, bold=True, clr=C_WHITE, al=PP_ALIGN.CENTER)

# ===== SLIDE 1 — TITLE =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl, C_DARK)
add_rect(sl, 0, 0, 13.33, 0.1, fill=C_ACCENT)
add_rect(sl, 0, 0.1, 13.33, 3.0, fill=C_NAVY)
add_rect(sl, 0, 3.1, 13.33, 0.06, fill=C_BLUE)
txt(sl, 'LedgerLM', 0.5, 0.28, 10, 1.05, sz=52, bold=True)
txt(sl, 'Multi-Tenant Platform Architecture', 0.5, 1.35, 12, 0.65, sz=25, clr=C_ACCENT)
txt(sl, 'Building Blocks  ·  Design Decisions  ·  Implementation Flow', 0.5, 2.05, 12, 0.5, sz=13, clr=C_MID, it=True)
for i,(lbl,desc,c) in enumerate([('Matasma / LedgerLM.ai','Platform owner — Super Admin control plane',C_ACCENT),
    ('Bosch','Tenant → bosch.ledgerlm.ai or self-hosted',C_BLUE),
    ('Nemko','Tenant → nemko.ledgerlm.ai',C_BLUE),
    ('Any New Client','Configured in Super Admin — minutes',C_GREEN)]):
    x=0.4+i*3.22
    add_rect(sl,x,3.32,3.08,0.92,fill=C_NAVY,line=c,lw=Pt(1.5))
    txt(sl,lbl,x+0.12,3.38,2.88,0.34,sz=10,bold=True,clr=c)
    txt(sl,desc,x+0.12,3.74,2.88,0.4,sz=8,clr=C_MID,it=True)
txt(sl,'10 Building Blocks  ·  April 2025  ·  Matasma Internal  ·  Confidential',0.4,6.95,12.5,0.38,sz=9,clr=C_MID,al=PP_ALIGN.CENTER,it=True)

# ===== SLIDE 2 — PLATFORM OVERVIEW =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Platform Overview','Three-layer model — Control Plane → App Layer → Data Layer')
layers=[
    ('CONTROL PLANE','ledgerlm.ai  (Matasma Super Admin)',C_ACCENT,
     ['Create & configure tenant domains','Set AI model, embedding, SSO, SMTP per domain',
      'Assign data_schema_type (Bosch / Nemko / Generic)','Monitor all tenants from one portal']),
    ('APPLICATION LAYER','Shared SaaS  or  Self-Hosted per tenant',C_BLUE,
     ['bosch.ledgerlm.ai → Bosch employees only','nemko.ledgerlm.ai → Nemko employees only',
      'Host header → middleware reads domain → loads config','Same Docker image for SaaS and self-hosted']),
    ('DATA LAYER','PostgreSQL + pgvector  (isolated by cube_id)',C_GREEN,
     ['Shared tables — all tenants isolated by cube_id','Vector embeddings per domain at correct dimension',
      'Row-level enforcement prevents cross-tenant access','SQL logic driven by data_schema_type from DB']),
]
for i,(title,sub,c,pts) in enumerate(layers):
    y=1.5+i*1.82
    add_rect(sl,0.35,y,12.62,1.65,fill=C_NAVY,line=c,lw=Pt(1.5))
    add_rect(sl,0.35,y,2.55,1.65,fill=c)
    txt(sl,title,0.42,y+0.1,2.4,0.45,sz=10,bold=True,clr=C_DARK,al=PP_ALIGN.CENTER)
    txt(sl,sub,0.42,y+0.62,2.4,0.72,sz=8,clr=C_DARK,al=PP_ALIGN.CENTER,it=True)
    for j,pt in enumerate(pts):
        txt(sl,f'• {pt}',3.05+(j%2)*4.85,y+0.2+(j//2)*0.72,4.7,0.65,sz=9,clr=C_LIGHT)
txt(sl,'→  The 10 building blocks on the following slides make all three layers work together.',0.35,7.1,12.62,0.3,sz=9.5,clr=C_ACCENT,it=True,al=PP_ALIGN.CENTER)

# ===== SLIDE 3 — BB1: TENANT IDENTITY =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 1 — Tenant Identity & Domain Creation','The "domains" table row is the entire tenant config — one record drives everything')
txt(sl,'Every other building block reads from this one DB record. Super Admin fills it once — nothing else needs changing.',0.4,1.42,12.5,0.36,sz=9.5,clr=C_MID,it=True)
cxs=[0.35,2.65,5.2,9.75]; cws=[2.25,2.5,4.5,1.35]; hdrs=['Column','Example Value','Purpose','Status']
hy=1.86
for ci,(h,cx,cw) in enumerate(zip(hdrs,cxs,cws)):
    add_rect(sl,cx,hy,cw,0.35,fill=C_BLUE); txt(sl,h,cx+0.07,hy+0.05,cw-0.14,0.25,sz=9,bold=True)
rows=[
    ('name','in.bosch.com','Domain identifier — matched from Host header','Done'),
    ('display_name','Bosch Financial Intel.','Shown in UI headers and emails','Done'),
    ('slug','bosch','URL prefix → bosch.ledgerlm.ai','Add'),
    ('status','active / suspended / trial','Tenant lifecycle management','Add'),
    ('plan_type','saas / self_hosted','Deployment mode determines data sovereignty','Add'),
    ('data_schema_type','bosch / nemko / generic','Drives Python SQL builder — KEY column','Add ★'),
    ('deployment_url','bosch.ledgerlm.ai','The live URL for this tenant','Add'),
    ('smtp_host / smtp_port','smtp.bosch.com  /  587','Per-domain email delivery config','Missing'),
    ('branding_name','Bosch Financial Intel.','Branded name on login page & emails','Missing'),
    ('branding_logo_url','https://cdn.../logo.png','Client logo everywhere in UI','Missing'),
    ('ai_embedding_dimensions','768 / 1536 / 3072','Vector size for this domain embeddings','Missing'),
]
for ri,row in enumerate(rows):
    ry=hy+0.35+ri*0.43
    add_rect(sl,0.35,ry,12.6,0.4,fill=C_NAVY if ri%2==0 else R(0x12,0x22,0x3A))
    st=row[3]; sc=C_GREEN if st=='Done' else (C_YELLOW if '★' in st else (C_ORANGE if st=='Add' else C_RED))
    for ci,(val,cx,cw) in enumerate(zip(row,cxs,cws)):
        c=sc if ci==3 else (C_ACCENT if ci==0 else C_LIGHT)
        txt(sl,val,cx+0.07,ry+0.08,cw-0.14,0.25,sz=8.5,bold=(ci==0),clr=c)

# ===== SLIDE 4 — BB2: AUTH =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 2 — Authentication Stack','Login method + SSO credentials + SMTP — each configured per domain independently')
cols=[
    ('2a  Login Method',C_ACCENT,['OTP (default) — 6-digit code via email','Microsoft SSO — Azure AD / Entra ID','Google SSO — OAuth2 client credentials','SAML — Enterprise IdP (future)','auth_method column in domains table','Login page adapts to configured method']),
    ('2b  Microsoft SSO\nConfig (per domain)',C_BLUE,['sso_tenant_id  →  Azure AD Tenant','sso_client_id  →  App Registration','sso_client_secret → encrypted in DB','Redirect URI auto-built from slug','One app handles multiple AD tenants','Schema: done  ·  Flow: partial']),
    ('2c  SMTP\n(per domain)',C_ORANGE,['smtp_host / smtp_port','smtp_user / smtp_password (encrypted)','from_address → noreply@bosch.ledgerlm.ai','from_name  →  Bosch LedgerLM','OTP + welcome emails from domain SMTP','STATUS: MISSING from schema entirely']),
    ('2d  Session / JWT',C_GREEN,['JWT scoped to domain_id + user_id','Bosch user cannot access Nemko data','Session expiry configurable per domain','Refresh token strategy per plan type','domain_id verified on every API call','Status: domain scoping done']),
]
for i,(title,c,pts) in enumerate(cols):
    x=0.35+i*3.25
    add_rect(sl,x,1.5,3.1,5.8,fill=C_NAVY,line=c,lw=Pt(1.5))
    add_rect(sl,x,1.5,3.1,0.62,fill=c)
    txt(sl,title,x+0.1,1.53,2.92,0.56,sz=10,bold=True,clr=C_DARK,al=PP_ALIGN.CENTER)
    for j,pt in enumerate(pts):
        tc=C_RED if 'MISSING' in pt or 'STATUS:' in pt else C_LIGHT
        txt(sl,f'• {pt}',x+0.13,2.22+j*0.77,2.88,0.7,sz=8.5,clr=tc)

# ===== SLIDE 5 — BB3: AI MODEL =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 3 — AI Model Configuration','Two independent configs per domain: Chat LLM + Embedding Model')
add_rect(sl,0.35,1.5,6.0,5.72,fill=C_NAVY,line=C_BLUE,lw=Pt(1.5))
add_rect(sl,0.35,1.5,6.0,0.55,fill=C_BLUE)
txt(sl,'3a  Chat / Reasoning LLM',0.45,1.54,5.82,0.46,sz=12,bold=True)
for i,(col,val,st) in enumerate([
    ('ai_provider','ollama  |  azure_openai  |  openai  |  anthropic','Done'),
    ('ai_endpoint','https://xxx.openai.azure.com  (Azure only)','Done'),
    ('ai_api_key','Encrypted at rest in DB','Done'),
    ('ai_chat_model','gpt-5.2-chat  or  qwen2.5:32b','Done'),
    ('ai_chat_api_version','2024-12-01-preview  (Azure only)','Done'),
    ('ai_system_prompt','Custom instructions / persona for this domain','Done'),
]):
    y=2.18+i*0.75
    add_rect(sl,0.5,y,5.7,0.66,fill=R(0x12,0x22,0x3A))
    txt(sl,col,0.62,y+0.05,2.1,0.28,sz=8.5,bold=True,clr=C_ACCENT)
    badge(sl,5.0,y+0.08,0.97,0.26,st,c=C_GREEN)
    txt(sl,val,0.62,y+0.35,4.3,0.27,sz=8,clr=C_LIGHT)
txt(sl,'✓ All 6 columns exist · Node.js reads + passes to streaming · Works today',0.45,6.62,5.75,0.45,sz=8,clr=C_GREEN,it=True)
add_rect(sl,6.65,1.5,6.3,5.72,fill=C_NAVY,line=C_ORANGE,lw=Pt(1.5))
add_rect(sl,6.65,1.5,6.3,0.55,fill=C_ORANGE)
txt(sl,'3b  Embedding Model',6.75,1.54,6.12,0.46,sz=12,bold=True)
for i,(col,val,st) in enumerate([
    ('ai_embedding_model','text-embedding-3-large  or  nomic-embed-text','Done'),
    ('ai_embedding_api_version','2024-02-01  (Azure only)','Done'),
    ('ai_embedding_dimensions','768  |  1536  |  3072','MISSING ★'),
]):
    y=2.18+i*0.88
    sc=C_GREEN if st=='Done' else C_RED
    add_rect(sl,6.8,y,6.0,0.76,fill=R(0x1A,0x25,0x3E),line=sc,lw=Pt(0.75))
    txt(sl,col,6.92,y+0.06,2.3,0.28,sz=8.5,bold=True,clr=C_ACCENT)
    badge(sl,11.6,y+0.1,1.2,0.28,st,c=sc)
    txt(sl,val,6.92,y+0.4,5.7,0.28,sz=8.5,clr=C_LIGHT)
add_rect(sl,6.65,4.85,6.3,2.22,fill=R(0x1A,0x10,0x05),line=C_ORANGE,lw=Pt(1))
txt(sl,'WHY DIMENSIONS MATTER',6.82,4.93,6.0,0.28,sz=9.5,bold=True,clr=C_ORANGE)
txt(sl,'pgvector stores vectors as vector(N) — a fixed N.\nBosch (Azure)  →  text-embedding-3-large  →  3072 dims\nNemko (Ollama) →  nomic-embed-text        →   768 dims\n\nIf Python uses wrong dims at ingest, RAG retrieval\nreturns garbage or crashes with dim mismatch error.\nThis column MUST exist and be read before any embed call.',6.82,5.25,6.0,1.72,sz=9,clr=C_LIGHT)

# ===== SLIDE 6 — BB4: VECTOR DB =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 4 — Vector DB Design (pgvector)','Dimension-awareness is non-negotiable for multi-tenant RAG')
add_rect(sl,0.35,1.5,12.62,0.66,fill=R(0x3D,0x10,0x10))
txt(sl,'Problem: pgvector requires vector(N) with fixed N at column creation. Azure embedding = 3072 dims. Ollama = 768 dims. Cannot mix in one column without a clear strategy.',0.52,1.58,12.22,0.5,sz=10,clr=C_LIGHT)
for i,(title,c,badge_txt,pts) in enumerate([
    ('Option A\nSingle column\n+ pad/truncate',C_RED,'AVOID',
     ['All stored at max 3072 dims','Smaller model padded with zeros','Padding distorts cosine similarity','Query accuracy degrades silently','Hard to debug — wrong results no error']),
    ('Option B\nTwo columns\n(Recommended)',C_GREEN,'USE THIS',
     ['embedding_768  vector(768)','embedding_3072 vector(3072)','Domain config picks correct column','No padding, no distortion','RAG query uses same column as ingest']),
    ('Option C\nSeparate tables\nper dim tier',C_BLUE,'VALID',
     ['document_chunks_768 table','document_chunks_3072 table','Cleanest isolation per provider','More tables but simpler columns','Query joins correct table per domain']),
]):
    x=0.35+i*4.35
    add_rect(sl,x,2.35,4.1,4.72,fill=C_NAVY,line=c,lw=Pt(2))
    add_rect(sl,x,2.35,4.1,0.72,fill=c)
    txt(sl,title,x+0.1,2.38,3.92,0.66,sz=10,bold=True,al=PP_ALIGN.CENTER)
    for j,pt in enumerate(pts): txt(sl,f'• {pt}',x+0.15,3.18+j*0.72,3.85,0.66,sz=9.5,clr=C_LIGHT)
    badge(sl,x+1.35,6.76,1.4,0.3,badge_txt,c=c)
add_rect(sl,0.35,7.1,12.62,0.3,fill=C_NAVY)
txt(sl,'Implementation: domains.ai_embedding_dimensions → Python reads → writes to embedding_768 or embedding_3072 → RAG uses same column → perfect similarity',0.5,7.15,12.2,0.22,sz=8.5,clr=C_ACCENT)

# ===== SLIDE 7 — BB5: DATA ISOLATION =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 5 — Data Schema & Tenant Isolation','Shared tables, cube_id as the isolation boundary')
txt(sl,'Ownership Hierarchy:',0.4,1.5,12.5,0.3,sz=11,bold=True,clr=C_ACCENT)
hier=[(0.5,'domains','LedgerLM\ntenant record',C_ACCENT),(3.5,'companies','Bosch GmbH\ncorp entity',C_BLUE),(6.5,'cubes','Data workspace\nJan–Feb 2025',C_BLUE),(9.5,'cube_fact_data','All rows filtered\nby cube_id',C_GREEN)]
for x,label,desc,c in hier:
    add_rect(sl,x,1.9,2.6,0.88,fill=c)
    txt(sl,label,x+0.1,1.93,2.42,0.36,sz=10,bold=True,al=PP_ALIGN.CENTER)
    txt(sl,desc,x+0.1,2.3,2.42,0.44,sz=8.5,al=PP_ALIGN.CENTER)
for i in range(3): txt(sl,'→',hier[i][0]+2.65,2.22,0.28,0.36,sz=14,bold=True,clr=C_MID,al=PP_ALIGN.CENTER)
txt(sl,'What keeps data safe today:',0.4,3.04,6.1,0.3,sz=10.5,bold=True,clr=C_ACCENT)
for i,p in enumerate(['Every query includes cube_id in WHERE clause','cube_id links to cubes.domain_id → authenticated domain','User login is scoped to their domain only','Python compile_sql receives cube_id from auth layer','Cube is created under a domain — no cross-domain assignment']):
    txt(sl,f'✓  {p}',0.5,3.45+i*0.56,6.1,0.5,sz=9.5,clr=C_LIGHT)
txt(sl,'Optional hardening (future):',6.8,3.04,6.1,0.3,sz=10.5,bold=True,clr=C_ORANGE)
for i,p in enumerate(['Postgres Row-Level Security (RLS) policy','Policy: user SELECTs only rows where cube_id IN (their cubes)','Enforced at DB level — code bugs cannot leak data','Audit log: who queried what + timestamp','Anomaly detection: flag cross-cube access attempts']):
    txt(sl,f'◎  {p}',6.9,3.45+i*0.56,6.1,0.5,sz=9.5,clr=C_MID)
add_rect(sl,0.35,6.45,12.62,0.84,fill=C_NAVY,line=C_BLUE,lw=Pt(1))
txt(sl,'Shared tables (all tenants):  cube_fact_data  ·  cube_plan_data  ·  document_chunks  ·  cubes  ·  domains  ·  users  ·  cube_query_jobs',0.5,6.53,12.2,0.3,sz=9.5,clr=C_LIGHT)
txt(sl,'Isolation boundary:  cube_id  (UUID → cubes.id → cubes.domain_id → must match authenticated user domain)',0.5,6.85,12.2,0.3,sz=9,clr=C_ACCENT,it=True)

# ===== SLIDE 8 — BB6: SQL ROUTING =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 6 — SQL Logic Routing  (data_schema_type)','Replace hardcoded name checks with a single DB config column')
txt(sl,'TODAY  (hardcoded)',0.35,1.5,6.0,0.34,sz=13,bold=True,clr=C_RED)
txt(sl,'AFTER  (DB-driven)',6.8,1.5,6.0,0.34,sz=13,bold=True,clr=C_GREEN)
add_rect(sl,6.65,1.45,0.06,5.88,fill=C_BLUE)
add_rect(sl,0.35,1.95,6.2,1.75,fill=C_DKRED)
for i,(ln,special) in enumerate([
    ("COMPANY_CONFIG = {",False),("  'bosch': { 'domains': ['in.bosch.com'] },",False),
    ("  'nemko': { 'domains': ['nemko.com'] }",False),("}", False),
    ("def is_nemko_domain(domain):  # name string check",True),
    ("    return 'nemko' in domain  # hardcoded!",True),
    ("nemko_pl_cube_id = '663678fc...'  # UUID x4!",True),
]):
    c=C_RED if special else RGBColor(0xFF,0xCC,0x70)
    txt(sl,ln,0.45,2.0+i*0.23,6.0,0.21,sz=8,clr=c)
txt(sl,'Problems:',0.35,3.8,6.1,0.28,sz=10,bold=True,clr=C_RED)
for i,p in enumerate(['New client requires Python code edit + redeploy','Nemko cube_id hardcoded UUID in 4 Python locations','Nemko-style data impossible for a new client without code change','Everything unknown defaults to Bosch — silent mis-routing']):
    txt(sl,f'✗  {p}',0.45,4.17+i*0.62,6.1,0.56,sz=9,clr=C_LIGHT)
add_rect(sl,6.75,1.95,6.18,1.3,fill=C_DKGRN)
for i,(ln,highlight) in enumerate([
    ("schema_type = domain_config['data_schema_type']",False),
    ("# 'bosch' | 'nemko' | 'generic'",False),
    ("cube_id = domain_config['cube_id']  # from DB",True),
    ("# Route to SQL builder by schema_type",False),
]):
    c=RGBColor(0x80,0xFF,0xA0) if highlight else RGBColor(0xA0,0xFF,0xC0)
    txt(sl,ln,6.85,2.0+i*0.3,6.0,0.27,sz=8,clr=c)
txt(sl,'Benefits:',6.75,3.38,6.0,0.28,sz=10,bold=True,clr=C_GREEN)
for i,p in enumerate(['New client → set data_schema_type in Super Admin → no code change','Cube ID always from DB — zero hardcoded UUIDs in Python','Another client with Bosch-style data: just set type = "bosch"','"generic" type → base SQL builder, no custom logic needed','Remove COMPANY_CONFIG dict entirely from codebase']):
    txt(sl,f'✓  {p}',6.85,3.75+i*0.63,6.0,0.56,sz=9,clr=C_LIGHT)

# ===== SLIDE 9 — BB7+8: BRANDING + DEPLOYMENT =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Blocks 7 & 8 — Branding + Deployment Modes','What clients see and how the app is delivered to them')
add_rect(sl,0.35,1.5,6.2,5.78,fill=C_NAVY,line=C_ACCENT,lw=Pt(1.5))
add_rect(sl,0.35,1.5,6.2,0.55,fill=C_ACCENT)
txt(sl,'Block 7  —  Per-Domain Branding',0.48,1.53,6.0,0.46,sz=12,bold=True,clr=C_DARK)
txt(sl,'New DB columns needed (all currently missing):',0.48,2.13,5.9,0.28,sz=9,clr=C_MID,it=True)
for i,(col,val,desc) in enumerate([
    ('branding_name','Bosch Financial Intelligence','Display name — login page + app header + emails'),
    ('branding_logo_url','https://cdn.../bosch-logo.png','Client logo — replaces LedgerLM logo everywhere'),
    ('branding_primary_color','#005691','Primary accent colour in UI (future)'),
    ('favicon_url','https://cdn.../favicon.ico','Browser tab icon (future)'),
]):
    y=2.5+i*0.72
    add_rect(sl,0.5,y,5.9,0.63,fill=R(0x12,0x22,0x3A))
    txt(sl,col,0.62,y+0.04,2.18,0.27,sz=8.5,bold=True,clr=C_ACCENT)
    txt(sl,val,0.62,y+0.33,2.18,0.25,sz=7.5,clr=C_MID,it=True)
    txt(sl,desc,2.88,y+0.15,3.4,0.28,sz=8.5,clr=C_LIGHT)
txt(sl,'Where branding appears:',0.48,5.44,5.9,0.28,sz=9,bold=True,clr=C_MID)
for i,p in enumerate(['Login page — logo + display name (instead of LedgerLM)','App header — client logo + name replaces platform logo','OTP email — from_name + logo in email body','Browser tab — "{client} | LedgerLM" title']):
    txt(sl,f'→ {p}',0.6,5.78+i*0.36,5.8,0.32,sz=9,clr=C_LIGHT)
add_rect(sl,6.78,1.5,6.2,5.78,fill=C_NAVY,line=C_BLUE,lw=Pt(1.5))
add_rect(sl,6.78,1.5,6.2,0.55,fill=C_BLUE)
txt(sl,'Block 8  —  Deployment Modes',6.9,1.53,6.0,0.46,sz=12,bold=True)
add_rect(sl,6.93,2.15,5.9,2.42,fill=R(0x0D,0x1D,0x35),line=C_ACCENT,lw=Pt(1))
txt(sl,'SaaS  (Shared Infrastructure)',7.05,2.2,5.7,0.32,sz=10.5,bold=True,clr=C_ACCENT)
for i,p in enumerate(['Nginx routes by subdomain → Host header → middleware sets domain','bosch.ledgerlm.ai and nemko.ledgerlm.ai on same server','All tenant data in Matasma DB — fastest to provision','plan_type = "saas"  ·  deployment_url column in domains']):
    txt(sl,f'• {p}',7.05,2.6+i*0.47,5.7,0.42,sz=8.5,clr=C_LIGHT)
add_rect(sl,6.93,4.72,5.9,2.42,fill=R(0x1A,0x10,0x05),line=C_ORANGE,lw=Pt(1))
txt(sl,'Self-Hosted  (Client Infrastructure)',7.05,4.77,5.7,0.32,sz=10.5,bold=True,clr=C_ORANGE)
for i,p in enumerate(['Docker image deployed on Bosch / Nemko own servers','Client controls their DB + Azure OpenAI — full data sovereignty','Data never leaves client infrastructure — regulatory compliance','plan_type = "self_hosted"  ·  deployment_url = client URL']):
    txt(sl,f'• {p}',7.05,5.18+i*0.47,5.7,0.42,sz=8.5,clr=C_LIGHT)

# ===== SLIDE 10 — BB9: SUPER ADMIN =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 9 — Super Admin Portal  (Matasma Control Plane)','One form at ledgerlm.ai to configure every tenant completely — no code changes needed')
sections=[
    ('Identity',C_ACCENT,['Domain name / display name','Slug (URL prefix)','Status — active/trial/suspended','Plan type + deployment URL']),
    ('Auth',C_BLUE,['Login method (OTP / SSO / Google)','SSO tenant_id + client_id','Session expiry config','Allowed email domains']),
    ('SMTP',C_ORANGE,['Host / port / credentials','from_address + from_name','Test email send button','STATUS: MISSING in schema']),
    ('AI Chat',C_GREEN,['Provider selection dropdown','Endpoint + encrypted API key','Chat model + API version','Custom system prompt textarea']),
    ('AI Embed',C_PURPLE,['Embedding model name','API version (Azure only)','Dimensions: 768/1536/3072','STATUS: dimensions MISSING']),
    ('Data',C_YELLOW,['data_schema_type picker','View assigned cubes','Trigger data re-ingestion','Ingestion history log']),
    ('Branding',C_PINK,['Display name','Logo URL + preview','Primary colour picker','Preview branded login page']),
    ('Deployment',C_MID,['Mode: SaaS / self-hosted','Deployment URL config','Health check status dot','Active users / last seen']),
]
for i,(title,c,pts) in enumerate(sections):
    ci=i%4; ri=i//4; x=0.35+ci*3.25; y=1.5+ri*2.88
    add_rect(sl,x,y,3.1,2.72,fill=C_NAVY,line=c,lw=Pt(1.5))
    add_rect(sl,x,y,3.1,0.5,fill=c)
    txt(sl,title,x+0.1,y+0.07,2.92,0.38,sz=11,bold=True,al=PP_ALIGN.CENTER,clr=C_WHITE)
    for j,pt in enumerate(pts):
        tc=C_RED if 'MISSING' in pt or 'STATUS:' in pt else C_LIGHT
        txt(sl,f'• {pt}',x+0.13,y+0.6+j*0.52,2.88,0.47,sz=8.5,clr=tc)

# ===== SLIDE 11 — BB10: INGESTION =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Block 10 — Data Ingestion Pipeline','Domain config must flow through every step — from parse to embed to store')
steps=[
    ('1  Auth &\nDomain',C_ACCENT,'User uploads file → authenticated → domain_id → load full domain config from DB',False),
    ('2  Data\nParsing',C_BLUE,'Parse Excel / Anaplan based on data_schema_type  (Bosch columns differ from Nemko columns)',False),
    ('3  Fact\nStorage',C_GREEN,'INSERT rows into cube_fact_data with cube_id from domain config — isolation guaranteed',False),
    ('4  Text\nChunking',C_YELLOW,'Extract text from uploaded document → split into overlapping chunks for RAG retrieval',False),
    ('5  Embedding\nGeneration',C_ORANGE,'Call get_embeddings(ai_config) — ai_config from domain record\n→ Azure (3072 dims) or Ollama (768 dims)  ·  CURRENTLY: dims hardcoded to 3072 in code',True),
    ('6  Vector\nStorage',C_PURPLE,'Write to correct pgvector column based on ai_embedding_dimensions\nCURRENTLY: single column only — hardcoded dims — breaks for new providers',True),
    ('7  Index\nUpdate',C_BLUE,'Update cube metadata → mark ingestion complete → clear domain config cache',False),
]
for i,(label,c,desc,issue) in enumerate(steps):
    y=1.52+i*0.72
    add_rect(sl,0.35,y,1.45,0.68,fill=c)
    txt(sl,label,0.38,y+0.03,1.4,0.62,sz=9.5,bold=True,al=PP_ALIGN.CENTER)
    add_rect(sl,1.83,y,11.12,0.68,fill=C_NAVY,line=c,lw=Pt(0.75))
    for li,ln in enumerate(desc.split('\n')):
        tc=C_ORANGE if issue and ('CURRENTLY' in ln or 'hardcoded' in ln) else C_LIGHT
        txt(sl,ln,1.93,y+0.07+li*0.27,11.0,0.25,sz=8.5,clr=tc)

# ===== SLIDE 12 — ONBOARDING FLOW =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'End-to-End Flow — New Client Onboarding','From "we want LedgerLM" to live AI queries — all driven by the single domain DB record')
flow=[
    ('Step 1','Super Admin\nFills Form','At ledgerlm.ai\nName, slug,\nauth method,\nAI config,\ndimensions,\nschema type,\nbranding',C_ACCENT),
    ('Step 2','Domain Row\nInserted','One domains\ntable record\nAll config here\nNo code deploy\nNo restart',C_BLUE),
    ('Step 3','Client Admin\nInvited','Email via domain\nSMTP config\nBranded logo\n+ name\nSSO link auto',C_GREEN),
    ('Step 4','Data\nUploaded','Excel / Anaplan\nParsed by\nschema_type\nEmbedded at\ncorrect dims',C_YELLOW),
    ('Step 5','Users\nLogin','slug.ledgerlm.ai\nBranded login\nOTP or SSO\nSession scoped\nto domain',C_ORANGE),
    ('Step 6','AI Queries\nLive','Client LLM chat\nRAG in client\nvector store\nSQL via type\nClient data only',C_ACCENT),
]
for i,(step,title,desc,c) in enumerate(flow):
    x=0.28+i*2.18
    if i<5: txt(sl,'→',x+2.14,2.78,0.28,0.36,sz=14,bold=True,clr=C_MID,al=PP_ALIGN.CENTER)
    add_rect(sl,x,1.5,2.1,0.5,fill=c)
    txt(sl,step,x+0.05,1.53,2.02,0.42,sz=9.5,bold=True,clr=C_DARK,al=PP_ALIGN.CENTER)
    add_rect(sl,x,2.0,2.1,0.82,fill=C_NAVY,line=c,lw=Pt(1.5))
    txt(sl,title,x+0.07,2.05,1.97,0.72,sz=9,bold=True,clr=c,al=PP_ALIGN.CENTER)
    add_rect(sl,x,2.82,2.1,3.0,fill=R(0x12,0x1E,0x30))
    txt(sl,desc,x+0.1,2.88,1.92,2.85,sz=8.5,clr=C_LIGHT)
add_rect(sl,0.28,5.95,12.76,1.42,fill=C_NAVY,line=C_BLUE,lw=Pt(1))
txt(sl,'What drives each step automatically from the domain config record:',0.44,6.02,12.4,0.28,sz=9.5,bold=True,clr=C_ACCENT)
for i,d in enumerate(['Step 1 → single DB INSERT','Step 2 → Host header middleware → domain lookup','Step 3 → smtp_host/port/from_address/from_name','Step 4 → data_schema_type + ai_embedding_dimensions','Step 5 → branding_name + branding_logo_url + auth_method','Step 6 → ai_provider + ai_chat_model + data_schema_type']):
    txt(sl,f'• {d}',0.44+(i%3)*4.28,6.36+(i//3)*0.42,4.1,0.38,sz=8.5,clr=C_LIGHT)

# ===== SLIDE 13 — STATUS =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Current Status — Done vs Needed','')
sm2={'Done':C_GREEN,'Partial':C_ORANGE,'Missing':C_RED}
hy=1.5
for ci,(h,cx,cw) in enumerate(zip(['Building Block','Status','Notes'],[0.35,6.28,7.72],[5.88,1.38,5.2])):
    add_rect(sl,cx,hy,cw,0.34,fill=C_BLUE)
    txt(sl,h,cx+0.07,hy+0.05,cw-0.14,0.24,sz=9,bold=True)
rows2=[
    ('Tenant identity — domains table','Partial','Missing data_schema_type, deployment_url, branding_*, smtp_* columns'),
    ('Authentication — OTP login','Done','Fully working'),
    ('Authentication — Microsoft SSO','Partial','Schema done, redirect flow partial'),
    ('Authentication — SMTP per domain','Missing','Single global SMTP — no per-domain config in schema'),
    ('AI Chat config per domain','Done','Schema + Node.js reads + streams correctly'),
    ('AI Embedding model per domain','Partial','Model stored; ai_embedding_dimensions column missing'),
    ('Vector DB dimension handling','Missing','Single column; Python hardcodes 3072 — breaks non-Azure domains'),
    ('Data isolation by cube_id','Done','Application-layer enforcement in all SQL builders'),
    ('SQL routing by data_schema_type','Missing','Hardcoded is_nemko_domain() name check — not DB driven'),
    ('Nemko cube_id from DB','Missing','UUID hardcoded in 4 Python locations — must be DB lookup'),
    ('Branding per domain','Missing','No schema columns — all pages show LedgerLM branding'),
    ('Super Admin — complete config form','Partial','Has SSO + AI chat; missing SMTP, dims, branding, deployment'),
    ('Deployment mode / URL config','Missing','No deployment_url column, no Host header middleware'),
    ('Ingestion — dimension-aware embeds','Partial','ai_config passed but dims hardcoded in vector_store.py:355'),
]
for i,(block,status,note) in enumerate(rows2):
    y=hy+0.34+i*0.36
    bg2=C_NAVY if i%2==0 else R(0x12,0x22,0x3A)
    add_rect(sl,0.35,y,5.88,0.33,fill=bg2); add_rect(sl,6.28,y,1.38,0.33,fill=sm2[status]); add_rect(sl,7.72,y,5.2,0.33,fill=bg2)
    txt(sl,block,0.43,y+0.06,5.7,0.22,sz=8.5,clr=C_LIGHT)
    txt(sl,status,6.33,y+0.06,1.25,0.22,sz=8,bold=True,al=PP_ALIGN.CENTER)
    txt(sl,note,7.8,y+0.06,5.08,0.22,sz=7.5,clr=C_MID,it=bool(note))
for j,(label,c) in enumerate([('Done',C_GREEN),('Partial',C_ORANGE),('Missing',C_RED)]):
    badge(sl,0.35+j*1.85,6.65,1.65,0.3,label,c=c)

# ===== SLIDE 14 — ROADMAP =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); header(sl,'Build Roadmap — Recommended Order','Each phase unlocks the next — start with schema, end with deployment')
for i,(title,c,dur,pts) in enumerate([
    ('Phase 1\nSchema\nFoundation',C_ACCENT,'1–2 days',['Add data_schema_type','Add ai_embedding_dimensions','Add branding_name + logo_url','Add deployment_url + plan_type','Add smtp_* columns (6 fields)','Update Super Admin form','npm run db:push']),
    ('Phase 2\nPython\nRouting',C_BLUE,'2–3 days',['Remove COMPANY_CONFIG dict','Replace is_nemko_domain() with DB','Fix 4 hardcoded nemko cube UUIDs','cube_id always from DB config','Cache invalidation on change','Test: Bosch + Nemko unchanged','']),
    ('Phase 3\nVector DB\nDimensions',C_GREEN,'1–2 days',['Add embedding_768 vector(768)','Add embedding_3072 vector(3072)','Read dims from domain config','Ingest → correct column','RAG → query correct column','Re-embed existing if dims changed','']),
    ('Phase 4\nSMTP &\nBranding',C_ORANGE,'2–3 days',['SMTP: OTP from domain config','Login page reads domain branding','App header shows client logo','Email templates use client name','Host header middleware','Public domain-config API','']),
    ('Phase 5\nDeployment\nConfig',C_YELLOW,'3–5 days',['Nginx subdomain routing','Docker image parameterisation','Health check per tenant','Self-hosted client config','Admin deployment status','Monitoring + uptime alerts','']),
]):
    x=0.3+i*2.6
    add_rect(sl,x,1.5,2.46,5.8,fill=C_NAVY,line=c,lw=Pt(1.5))
    add_rect(sl,x,1.5,2.46,0.7,fill=c)
    txt(sl,title,x+0.08,1.52,2.3,0.64,sz=9.5,bold=True,al=PP_ALIGN.CENTER)
    badge(sl,x+0.73,2.27,1.0,0.28,dur,c=C_DARK)
    for j,pt in enumerate(pts):
        if pt: txt(sl,f'• {pt}',x+0.12,2.65+j*0.73,2.28,0.66,sz=8,clr=C_LIGHT)

prs.save('/home/runner/workspace/LedgerLM_Architecture.pptx')
print('SUCCESS')
